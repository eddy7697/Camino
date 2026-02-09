"""
朝聖之路 Camino de Santiago PowerPoint 簡報產生器
所有圖片皆依原始比例放置，不會拉伸變形。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "img")

# ── Color Palette ──
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_LIGHT = RGBColor(0xE8, 0xD5, 0x90)
CREAM = RGBColor(0xF8, 0xF4, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TERRACOTTA = RGBColor(0xB5, 0x56, 0x1A)
TEXT_DARK = RGBColor(0x3A, 0x35, 0x30)
TEXT_LIGHT = RGBColor(0x6B, 0x61, 0x58)
SUBTLE = RGBColor(0xBB, 0xBB, 0xCC)
DIM = RGBColor(0xAA, 0xAA, 0xBB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width   # in EMU
SLIDE_H = prs.slide_height


# ══════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════

def img_path(name):
    return os.path.join(IMG, name)


def get_ratio(name):
    """Return (width, height, w/h ratio) of an image file."""
    with Image.open(img_path(name)) as im:
        w, h = im.size
    return w, h, w / h


def fit_contain(img_name, max_w_inches, max_h_inches):
    """Calculate (w, h) in Inches that fits inside the box while keeping aspect ratio."""
    _, _, ratio = get_ratio(img_name)
    max_w = max_w_inches
    max_h = max_h_inches
    # Try fitting to width
    w = max_w
    h = max_w / ratio
    if h > max_h:
        # Fit to height instead
        h = max_h
        w = max_h * ratio
    return w, h


def add_img_contain(slide, img_name, box_left, box_top, box_w, box_h):
    """Add image centered inside a bounding box, maintaining aspect ratio."""
    fit_w, fit_h = fit_contain(img_name, box_w, box_h)
    # Center inside the box
    left = box_left + (box_w - fit_w) / 2
    top = box_top + (box_h - fit_h) / 2
    return slide.shapes.add_picture(
        img_path(img_name),
        Inches(left), Inches(top),
        Inches(fit_w), Inches(fit_h)
    )


def add_image_bg_cover(slide, img_name):
    """Add image as slide background using 'cover' strategy (may overflow)."""
    _, _, ratio = get_ratio(img_name)
    slide_w = 13.333
    slide_h = 7.5
    slide_ratio = slide_w / slide_h  # ~1.778

    if ratio >= slide_ratio:
        # Image is wider → fit to height, overflow width
        h = slide_h
        w = slide_h * ratio
    else:
        # Image is taller → fit to width, overflow height
        w = slide_w
        h = slide_w / ratio

    left = (slide_w - w) / 2
    top = (slide_h - h) / 2
    slide.shapes.add_picture(
        img_path(img_name),
        Inches(left), Inches(top),
        Inches(w), Inches(h)
    )


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_overlay(slide, alpha=0.55):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    from pptx.oxml.ns import qn
    from lxml import etree
    sp_pr = shape._element.spPr
    solid_fill = sp_pr.find(qn('a:solidFill'))
    if solid_fill is not None:
        clr_elem = solid_fill[0]
        alpha_elem = etree.SubElement(clr_elem, qn('a:alpha'))
        alpha_elem.set('val', str(int((1 - alpha) * 100000)))
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft JhengHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=18, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Microsoft JhengHei",
             space_before=Pt(6), space_after=Pt(6)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Pt(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title (Hero) — bg: 18.jpg (portrait 0.46)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_image_bg_cover(slide, "18.jpg")
add_overlay(slide, alpha=0.6)

add_textbox(slide, 1, 1.5, 11.3, 1.5,
            "朝聖之路", 60, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.0, 11.3, 0.7,
            "CAMINO DE SANTIAGO", 24, GOLD_LIGHT, False, PP_ALIGN.CENTER)
add_gold_line(slide, 5.5, 3.8, 2.3)
add_textbox(slide, 2.5, 4.2, 8.3, 1,
            "一段徒步穿越西班牙的信仰旅程，用雙腳丈量 800 公里的恩典之路",
            18, WHITE, False, PP_ALIGN.CENTER)

stats_y = 5.5
for i, (num, label) in enumerate([("32", "天"), ("800", "公里"), ("6", "同行者")]):
    x = 3.5 + i * 2.2
    add_textbox(slide, x, stats_y, 1.8, 0.8, num, 44, GOLD, True, PP_ALIGN.CENTER)
    add_textbox(slide, x, stats_y + 0.75, 1.8, 0.4, label, 14, WHITE, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: 踏上朝聖之路 — 01.jpg (portrait 0.75)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

# 01.jpg is portrait 3:4 → in box 4.2w x 5.9h → fit_contain → 4.2 x 5.6
add_img_contain(slide, "01.jpg", 0.8, 0.8, 4.2, 5.9)

add_textbox(slide, 5.8, 1.2, 6.5, 0.8, "踏上朝聖之路", 36, NAVY, True)
add_gold_line(slide, 5.8, 2.1, 1.5)

txBox = add_textbox(slide, 5.8, 2.5, 6.8, 2.5,
    "2025 年 5 月 10 日，背起行囊，從台灣出發前往法國巴黎蒙帕納斯，正式化身為背包客，踏上這段一生一次的朝聖旅程。",
    16, TEXT_DARK)
add_para(txBox.text_frame,
    "朝聖之路（Camino de Santiago）是一條跨越千年的信仰之路，從法國南部翻越庇里牛斯山，一路徒步穿越西班牙北部，最終抵達聖地牙哥德孔波斯特拉主座教堂。",
    16, TEXT_LIGHT, space_before=Pt(12))

quote_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.8), Inches(5.3), Inches(6.8), Inches(1.2))
quote_shape.fill.solid()
quote_shape.fill.fore_color.rgb = RGBColor(0xEF, 0xE8, 0xD8)
quote_shape.line.fill.background()
add_textbox(slide, 6.1, 5.45, 6.2, 0.9,
    "「走了 32 天的路，800 公里的信仰之旅，每一步都是恩典。」",
    17, TERRACOTTA)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: 巴黎 — 04.jpg (portrait 0.75)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, NAVY)

add_textbox(slide, 0.8, 0.5, 3, 0.4, "巴黎 PARIS", 13, GOLD, True)
add_textbox(slide, 0.8, 1.0, 5, 1,
            "在前往 SJPP 之前\n先與艾菲爾鐵塔合影", 30, WHITE, True)
add_gold_line(slide, 0.8, 2.5, 1.5)
add_textbox(slide, 0.8, 2.9, 5, 2.5,
            "從蒙帕納斯出發前往朝聖之路的起點 Saint-Jean-Pied-de-Port，途中抽空走到艾菲爾鐵塔，為這趟旅程留下浪漫的序章。",
            16, SUBTLE)

# 04.jpg portrait 0.75 → box 5.0w x 6.5h → contain → 4.875 x 6.5
add_img_contain(slide, "04.jpg", 7.2, 0.5, 5.3, 6.5)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: 星星鎮 — 05.jpg (portrait 0.75) + 06.jpg (landscape 1.33)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

add_textbox(slide, 0.8, 0.5, 3, 0.4, "西班牙", 13, TERRACOTTA, True)
add_textbox(slide, 0.8, 1.0, 11, 0.8, "星星鎮 Estella", 36, NAVY, True)
add_gold_line(slide, 0.8, 1.9, 1.5)
add_textbox(slide, 0.8, 2.3, 11.5, 0.8,
    "翻越庇里牛斯山後進入西班牙，沿途經過充滿中世紀風情的星星鎮。古老的石板路、溫暖的陽光，每一步都踏在歷史的印記上。",
    16, TEXT_LIGHT)

# 05.jpg portrait 3:4, 06.jpg landscape 4:3 — give each a 5.8w x 4.0h box
add_img_contain(slide, "05.jpg", 0.8, 3.3, 5.0, 4.0)
add_img_contain(slide, "06.jpg", 6.5, 3.3, 6.0, 4.0)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: Logroño — 07.jpg (portrait 0.75)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, NAVY)

# Photo on left — portrait
add_img_contain(slide, "07.jpg", 0.5, 0.5, 5.2, 6.5)

add_textbox(slide, 6.5, 0.5, 3, 0.4, "LOGROÑO", 13, GOLD, True)
add_textbox(slide, 6.5, 1.0, 6.3, 1.2, "聖瑪利亞主教座堂", 32, WHITE, True)
add_gold_line(slide, 6.5, 2.3, 1.5)
add_textbox(slide, 6.5, 2.7, 6.3, 3,
    "Logroño 最著名的景點 Concatedral de Santa María de la Redonda 主教座堂，歷史可以追溯到 15 世紀，教堂內珍藏米開朗基羅的油畫，並開放給朝聖者參觀。",
    16, SUBTLE)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: 朝聖者紀念碑 — 22-25.jpg (3 portrait + 1 landscape)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

add_textbox(slide, 0.8, 0.4, 4, 0.4, "聖羅克高地", 13, TERRACOTTA, True)
add_textbox(slide, 0.8, 0.9, 11, 0.8, "朝聖者紀念碑", 36, NAVY, True)
add_gold_line(slide, 0.8, 1.8, 1.5)
add_textbox(slide, 0.8, 2.1, 11.5, 1.0,
    "聖羅克高地上的朝聖者紀念碑，傳說這位朝聖者原是個惡霸流氓，如今卻成為朝聖之路上最重要的象徵之一。許多朝聖者會在他的腳上貼上 OK 繃——因為走了這麼遠的路，誰的腳不起水泡呢？",
    15, TEXT_LIGHT)

# 22(P), 23(P), 24(P), 25(L) — each in a ~3.0 x 3.8 box
imgs_row = ["22.jpg", "23.jpg", "24.jpg", "25.jpg"]
box_w = 2.9
box_h = 3.8
gap = 0.25
start_x = 0.8
for i, fname in enumerate(imgs_row):
    x = start_x + i * (box_w + gap)
    add_img_contain(slide, fname, x, 3.3, box_w, box_h)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: 美食 — 08.jpg (portrait 0.56) + 10.jpg (portrait 0.75)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, NAVY)

add_textbox(slide, 0.8, 0.3, 11.7, 0.4, "朝聖路上的美食", 13, GOLD, True, PP_ALIGN.CENTER)

# 08.jpg very tall portrait (0.56) → box 5.5w x 5.0h
add_img_contain(slide, "08.jpg", 0.8, 1.2, 5.5, 5.0)
add_textbox(slide, 0.8, 6.3, 5.5, 0.5,
    "心心念念的蒜蘑菇", 18, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 6.8, 5.5, 0.5,
    "在台灣就心心念念的西班牙蒜蘑菇，終於品嚐到了！",
    12, DIM, False, PP_ALIGN.CENTER)

# 10.jpg portrait (0.75) → box 5.5w x 5.0h
add_img_contain(slide, "10.jpg", 7.0, 1.2, 5.5, 5.0)
add_textbox(slide, 7.0, 6.3, 5.5, 0.5,
    "薩里亞的水煮章魚", 18, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 7.0, 6.8, 5.5, 0.5,
    "進入 Sarria 前的音樂 Bar，好吃！",
    12, DIM, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: 朝聖者護照 — 11.jpg (landscape 16:9)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

add_textbox(slide, 0.8, 0.5, 4, 0.4, "朝聖印記", 13, TERRACOTTA, True)
add_textbox(slide, 0.8, 1.0, 5.5, 0.8, "朝聖者護照", 36, NAVY, True)
add_gold_line(slide, 0.8, 2.0, 1.5)
add_textbox(slide, 0.8, 2.4, 5.5, 3,
    "每經過一個小鎮就可獲得一個紀念章，抵達目的地時，朝聖者也依此獲發朝聖證明。\n\n早上 10:30，同行 6 人第一個到達倒數 100 公里處，一口氣走了近 15 公里！",
    16, TEXT_LIGHT)

# 11.jpg is landscape 16:9 → box 6.0w x 6.0h → will be wide
add_img_contain(slide, "11.jpg", 6.5, 0.8, 6.2, 6.2)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: 倒數 100km — 12.jpg + 13.jpg (both portrait 0.75)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, NAVY)

add_textbox(slide, 0.8, 0.3, 3, 0.4, "100 KM", 13, GOLD, True)
add_textbox(slide, 0.8, 0.8, 11.5, 0.8,
    "再踏出一步就破百了！", 36, WHITE, True, PP_ALIGN.CENTER)
add_gold_line(slide, 5.8, 1.8, 1.8)
add_textbox(slide, 2, 2.1, 9.3, 0.7,
    "800 公里的路，已經走了 700 公里，終點就在前方。這一刻的激動難以言喻。",
    16, SUBTLE, False, PP_ALIGN.CENTER)

# Both portrait → each in a 5.5w x 4.3h box
add_img_contain(slide, "12.jpg", 0.8, 3.0, 5.5, 4.3)
add_img_contain(slide, "13.jpg", 7.0, 3.0, 5.5, 4.3)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: 終點教堂 — bg: 18.jpg (portrait 0.46)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_image_bg_cover(slide, "18.jpg")
add_overlay(slide, alpha=0.5)

add_textbox(slide, 0.8, 0.5, 4, 0.4, "終點 SANTIAGO", 13, GOLD, True)
add_textbox(slide, 1, 2.0, 11.3, 1.5,
    "甩帽畢業了！", 52, WHITE, True, PP_ALIGN.CENTER)
add_gold_line(slide, 5.5, 3.8, 2.3)
add_textbox(slide, 2, 4.3, 9.3, 2,
    "朝聖之路的終點——聖地牙哥德孔波斯特拉主座教堂\n在雨中抵達這座宏偉的教堂前，將帽子拋向天空\n32 天的堅持與信念，在這一刻化為最美的回憶",
    18, WHITE, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: 朝聖者證書 — 26.jpg (portrait 0.75)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

# Portrait photo on left
add_img_contain(slide, "26.jpg", 0.8, 0.8, 4.5, 5.9)

add_textbox(slide, 6.0, 0.5, 4, 0.4, "榮耀時刻", 13, TERRACOTTA, True)
add_textbox(slide, 6.0, 1.2, 6.5, 1, "拿到朝聖者證書了！", 34, NAVY, True)
add_gold_line(slide, 6.0, 2.3, 1.5)
add_textbox(slide, 6.0, 2.7, 6.5, 3,
    "走了 32 天的路，終於拿到朝聖者證書了。\n\n走了 800 公里到聖地牙哥-德孔波斯特拉教堂的那一刻——快哭了。\n\n手中的兩張證書，是信仰與毅力的最佳見證。",
    16, TEXT_LIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: 世界盡頭 — bg: 16.jpg (landscape 1.33, 有海景)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_image_bg_cover(slide, "16.jpg")
add_overlay(slide, alpha=0.55)

add_textbox(slide, 1, 1.2, 11.3, 1.2,
    "世界的盡頭", 48, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 1, 2.5, 11.3, 0.5,
    "FINISTERRE", 20, GOLD_LIGHT, False, PP_ALIGN.CENTER)
add_gold_line(slide, 5.5, 3.3, 2.3)

txBox = add_textbox(slide, 2.5, 3.8, 8.3, 3,
    "6 月 14 日中午 11:50，來到了菲斯特雷角加利西亞海岸",
    17, WHITE, False, PP_ALIGN.CENTER)
add_para(txBox.text_frame, "歸零里程碑 Km 0,000",
         17, GOLD_LIGHT, False, PP_ALIGN.CENTER, space_before=Pt(10))
add_para(txBox.text_frame, "象徵著一切歸零，從頭開始",
         17, WHITE, False, PP_ALIGN.CENTER, space_before=Pt(10))
add_para(txBox.text_frame, "願倒空自己，讓心歸零",
         22, GOLD, True, PP_ALIGN.CENTER, space_before=Pt(16))


# ══════════════════════════════════════════════════════════════
# SLIDE 13: Finisterre 照片 — 15(P) + 16(L) + 17(P)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, NAVY)

# 15.jpg portrait, 16.jpg landscape, 17.jpg portrait
# Give each a box, respecting orientation
# Layout: 3 columns, each ~4.0w x 6.5h
add_img_contain(slide, "15.jpg", 0.4, 0.5, 3.9, 6.5)
add_img_contain(slide, "16.jpg", 4.5, 0.5, 4.4, 6.5)
add_img_contain(slide, "17.jpg", 9.1, 0.5, 3.9, 6.5)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: 羅卡角 — 19.jpg (portrait 0.75) + 20.jpg (landscape 1.33)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

add_textbox(slide, 0.8, 0.5, 4, 0.4, "CABO DA ROCA", 13, TERRACOTTA, True)
add_textbox(slide, 0.8, 1.0, 11.5, 1,
    "陸止於此、海始於斯", 38, NAVY, True, PP_ALIGN.CENTER)
add_gold_line(slide, 5.5, 2.1, 2.3)
add_textbox(slide, 2, 2.5, 9.3, 0.7,
    "6 月 19 日來到葡萄牙 Roca 羅卡角——歐洲大陸的最西端，有人稱這也是另一個世界的盡頭。",
    16, TEXT_LIGHT, False, PP_ALIGN.CENTER)

# 19.jpg portrait → 20.jpg landscape
add_img_contain(slide, "19.jpg", 0.8, 3.5, 5.0, 3.8)
add_img_contain(slide, "20.jpg", 6.5, 3.5, 6.0, 3.8)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: 葡萄牙中心點 + 搭錯車 — 21.jpg (P) + 30.jpg (L)
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, NAVY)

# 21.jpg portrait → box 5.5w x 4.8h
add_img_contain(slide, "21.jpg", 0.5, 0.8, 5.5, 4.8)
add_textbox(slide, 0.5, 5.8, 5.5, 0.5,
    "葡萄牙最中心點", 20, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 6.3, 5.5, 0.5,
    "從這個點可以到達葡萄牙的每個城市",
    13, DIM, False, PP_ALIGN.CENTER)

# 30.jpg landscape → box 6.3w x 4.8h
add_img_contain(slide, "30.jpg", 6.5, 0.8, 6.3, 4.8)
add_textbox(slide, 6.5, 5.8, 6.3, 0.5,
    "沒錯！我們搭錯車了", 20, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, 6.5, 6.3, 6.3, 0.5,
    "旅途中的小插曲，也成了最難忘的回憶",
    13, DIM, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 16: 旅途光影 (1/2) — 4 images
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

add_textbox(slide, 0.8, 0.2, 11.5, 0.8,
    "旅途光影", 34, NAVY, True, PP_ALIGN.CENTER)
add_gold_line(slide, 5.8, 1.1, 1.8)
add_textbox(slide, 2, 1.3, 9.3, 0.5,
    "沿途記錄下的美好瞬間", 14, TEXT_LIGHT, False, PP_ALIGN.CENTER)

gallery_page1 = ["02.jpg", "03.jpg", "09.jpg", "14.jpg"]
box_w = 2.8
box_h = 5.0
gap = 0.3
total_w = 4 * box_w + 3 * gap
start_x = (13.333 - total_w) / 2
for i, fname in enumerate(gallery_page1):
    x = start_x + i * (box_w + gap)
    add_img_contain(slide, fname, x, 2.1, box_w, box_h)


# ══════════════════════════════════════════════════════════════
# SLIDE 17: 旅途光影 (2/2) — 3 images
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, CREAM)

add_textbox(slide, 0.8, 0.2, 11.5, 0.8,
    "旅途光影", 34, NAVY, True, PP_ALIGN.CENTER)
add_gold_line(slide, 5.8, 1.1, 1.8)

gallery_page2 = ["27.jpg", "28.jpg", "29.jpg"]
box_w2 = 3.2
box_h2 = 5.3
gap2 = 0.4
total_w2 = 3 * box_w2 + 2 * gap2
start_x2 = (13.333 - total_w2) / 2
for i, fname in enumerate(gallery_page2):
    x = start_x2 + i * (box_w2 + gap2)
    add_img_contain(slide, fname, x, 1.6, box_w2, box_h2)


# ══════════════════════════════════════════════════════════════
# SLIDE 18: 感恩結語
# ══════════════════════════════════════════════════════════════
slide = new_slide()
add_bg(slide, NAVY)

add_textbox(slide, 1, 0.8, 11.3, 1,
    "感恩 · 歸零 · 再出發", 42, GOLD, True, PP_ALIGN.CENTER)
add_gold_line(slide, 5.5, 2.0, 2.3)

txBox = add_textbox(slide, 2, 2.5, 9.3, 1.5,
    "32 天，800 公里，從法國巴黎到西班牙聖地牙哥，再到世界的盡頭菲斯特雷角。",
    17, WHITE, False, PP_ALIGN.CENTER)
add_para(txBox.text_frame,
    "每一步都是信心的操練，每一天都是恩典的經歷。",
    17, WHITE, False, PP_ALIGN.CENTER, space_before=Pt(8))
add_para(txBox.text_frame,
    "這不只是一段徒步旅行，更是一場與自己、與信仰的深度對話。",
    17, WHITE, False, PP_ALIGN.CENTER, space_before=Pt(8))

verse_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3), Inches(4.8), Inches(7.3), Inches(1.5))
verse_shape.fill.solid()
verse_shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x3A)
verse_shape.line.color.rgb = RGBColor(0x40, 0x3A, 0x28)
verse_shape.line.width = Pt(1)

add_textbox(slide, 3.3, 5.0, 6.7, 0.6,
    "「你的話是我腳前的燈，是我路上的光。」",
    20, GOLD_LIGHT, False, PP_ALIGN.CENTER)
add_textbox(slide, 3.3, 5.65, 6.7, 0.4,
    "—— 詩篇 119:105",
    14, RGBColor(0x99, 0x99, 0xAA), False, PP_ALIGN.CENTER)
add_textbox(slide, 2, 6.5, 9.3, 0.6,
    "願將這段旅程的感動，與教會的弟兄姊妹們分享",
    17, GOLD_LIGHT, False, PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# Slide Transitions
# ══════════════════════════════════════════════════════════════
from pptx.oxml.ns import qn
from lxml import etree


def add_transition(slide, trans_type="fade", speed="med", advance_ms=None):
    """Add transition effect to a slide via XML.
    trans_type: fade, push, wipe, cover, split, blinds, dissolve
    speed: slow, med, fast
    advance_ms: auto-advance after N milliseconds (None = click only)
    """
    sld = slide._element
    # Remove existing transition
    for old in sld.findall(qn('p:transition')):
        sld.remove(old)

    trans = etree.SubElement(sld, qn('p:transition'))
    trans.set('spd', speed)
    trans.set('advClick', '1')
    if advance_ms is not None:
        trans.set('advTm', str(advance_ms))

    if trans_type == "fade":
        etree.SubElement(trans, qn('p:fade'))
    elif trans_type == "push":
        child = etree.SubElement(trans, qn('p:push'))
        child.set('dir', 'l')
    elif trans_type == "wipe":
        child = etree.SubElement(trans, qn('p:wipe'))
        child.set('dir', 'd')
    elif trans_type == "cover":
        child = etree.SubElement(trans, qn('p:cover'))
        child.set('dir', 'l')
    elif trans_type == "split":
        child = etree.SubElement(trans, qn('p:split'))
        child.set('orient', 'horz')
        child.set('dir', 'out')
    elif trans_type == "blinds":
        child = etree.SubElement(trans, qn('p:blinds'))
        child.set('dir', 'vert')
    elif trans_type == "dissolve":
        etree.SubElement(trans, qn('p:dissolve'))


# Define transition per slide: (type, speed)
transitions = [
    ("fade",    "slow"),   # 1  Hero 封面
    ("fade",    "med"),    # 2  踏上朝聖之路
    ("push",    "med"),    # 3  巴黎
    ("fade",    "med"),    # 4  星星鎮
    ("push",    "med"),    # 5  Logroño
    ("fade",    "med"),    # 6  朝聖者紀念碑
    ("push",    "med"),    # 7  美食
    ("fade",    "med"),    # 8  朝聖者護照
    ("push",    "med"),    # 9  倒數 100km
    ("fade",    "slow"),   # 10 終點教堂
    ("fade",    "med"),    # 11 朝聖者證書
    ("dissolve","slow"),   # 12 世界盡頭
    ("fade",    "med"),    # 13 Finisterre 照片
    ("push",    "med"),    # 14 羅卡角
    ("fade",    "med"),    # 15 葡萄牙 + 搭錯車
    ("fade",    "med"),    # 16 旅途光影 1/2
    ("fade",    "med"),    # 17 旅途光影 2/2
    ("fade",    "slow"),   # 18 感恩結語
]

for i, slide in enumerate(prs.slides):
    if i < len(transitions):
        t_type, t_speed = transitions[i]
        add_transition(slide, t_type, t_speed)

print(f"Added transitions to {len(prs.slides)} slides.")


# ══════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════
output_path = os.path.join(BASE, "朝聖之路.pptx")
prs.save(output_path)
print(f"PowerPoint saved to: {output_path}")
